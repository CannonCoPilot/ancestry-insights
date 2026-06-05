"""Build executive summary PowerPoint slide deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

PROJECT = Path(__file__).resolve().parent.parent
OUT = PROJECT / "outputs"
DECK_PATH = PROJECT / "outputs" / "FamilySearch_User_Persistence_Analysis.pptx"

# Colors
GREEN_DARK = RGBColor(0x1a, 0x43, 0x14)
GREEN_MID = RGBColor(0x3b, 0x85, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_slide():
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_bg(slide, color=LIGHT_BG):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=14, bold=False,
             color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_image(slide, path, left, top, width=None, height=None):
    """Add an image if it exists."""
    p = Path(path)
    if not p.exists():
        print(f"  WARNING: {p} not found")
        return None
    kwargs = {}
    if width: kwargs["width"] = Inches(width)
    if height: kwargs["height"] = Inches(height)
    if not kwargs:
        kwargs["width"] = Inches(6)
    return slide.shapes.add_picture(str(p), Inches(left), Inches(top), **kwargs)


def title_bar(slide, title, subtitle=None):
    """Add a dark green title bar across the top."""
    # Green bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.2))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_DARK
    shape.line.fill.background()

    add_text(slide, 0.5, 0.15, 12, 0.6, title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.5, 0.7, 12, 0.4, subtitle, font_size=14, color=RGBColor(0xC0, 0xE0, 0xB0))


def caption_box(slide, left, top, width, text):
    """Add a caption below a figure."""
    add_text(slide, left, top, width, 1.0, text, font_size=11, color=GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s, GREEN_DARK)
add_text(s, 1.5, 1.5, 10, 1.5,
         "Analysis of FamilySearch\nUser Persistence",
         font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.5, 10, 0.8,
         "What Predicts Whether Users Stay Engaged with Family History?",
         font_size=22, color=RGBColor(0xA0, 0xD8, 0x7A), alignment=PP_ALIGN.CENTER)
add_text(s, 1.5, 5.0, 10, 0.5,
         "Hypothesis-Driven Discriminant Analysis  |  7.6M User Accounts  |  30 Behavioral + Contextual Features",
         font_size=14, color=RGBColor(0x80, 0xB0, 0x70), alignment=PP_ALIGN.CENTER)
add_text(s, 1.5, 6.2, 10, 0.5,
         "March 2026",
         font_size=14, color=RGBColor(0x80, 0xB0, 0x70), alignment=PP_ALIGN.CENTER)
print("Slide 1: Title")

# ═══════════════════════════════════════════════════
# SLIDE 2: The Research Question
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "The Research Question")
add_text(s, 0.8, 1.5, 11.5, 1.5,
         '"Should we expect that people of different cultures, ages, education, or wealth '
         'will differ greatly in what captivates and holds their attention? '
         'Is engagement with the story of our family histories more influenced by our '
         'circumstances — or by how we approach that engagement?"',
         font_size=20, color=BLACK, alignment=PP_ALIGN.CENTER)
add_text(s, 0.8, 3.5, 11.5, 2.5,
         "We test two competing hypotheses:\n\n"
         "H1 (Engagement-Driven): User Persistence is predicted by behavioral engagement patterns — "
         "specifically Velocity (how quickly users progress through milestones), Volume (rate of contributions), "
         "and Sequencing (breadth and order of activities).\n\n"
         "H0 (Context-Driven): User Persistence is predicted by demographic and contextual factors — "
         "age, country of origin, economic development, religiosity, and digital infrastructure.",
         font_size=15, color=BLACK)
print("Slide 2: Research Question")

# ═══════════════════════════════════════════════════
# SLIDE 3: The Dataset — Column Density
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "The Dataset: 7.6M Users, but Sparse Beyond Core Activities")
add_image(s, OUT / "data-density/fig1_column_density.png", 0.3, 1.4, width=7.5)
add_text(s, 8.0, 1.5, 4.8, 5.5,
         "33 columns across demographics, activity counts, milestone dates, and name contributions.\n\n"
         "Key finding: Only 7 columns have >50% meaningful (non-null, non-zero) data. "
         "The 10.2% MNAR block (red) affects all 11 activity columns identically — a pipeline artifact, not random missingness.\n\n"
         "Logins, tree edits, and names are the three activities with enough prevalence to support rate-based features. "
         "Sources (8%), memories (2.3%), and record edits (0.6%) are too sparse for individual modeling but contribute to activity breadth counts.\n\n"
         "Province and City are 97% unknown — but the 3% that exist are exclusively Member accounts.",
         font_size=12, color=BLACK)
print("Slide 3: Column Density")

# ═══════════════════════════════════════════════════
# SLIDE 4: Population Funnel
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "From 7.6M Accounts to 3.2M Analytically Viable Users")
add_image(s, OUT / "data-density/fig3_density_funnel.png", 0.3, 1.4, width=5.5)
add_image(s, OUT / "phase4/fig_population_funnel.png", 6.5, 1.4, width=6.3)
caption_box(s, 0.3, 5.8, 5.5,
            "Left: Progressive attrition by data density tier. "
            "Right: The Phase 4 pipeline funnel from raw data to analytical subsamples.")
add_text(s, 0.3, 6.5, 12.5, 0.8,
         "42% of users (Tier D) have login + tree edits + names + milestone dates — sufficient for Volume, Velocity, "
         "and Sequencing constructs. 10.2% are untracked (MNAR exclusion); 37% are single-browse visitors (Segment 0).",
         font_size=12, color=GRAY)
print("Slide 4: Funnel")

# ═══════════════════════════════════════════════════
# SLIDE 5: External Enrichment
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Country-Level Enrichment: 26 Variables Across 6 Global Data Sources")
add_image(s, OUT / "phase3/fig_source_coverage.png", 0.3, 1.4, width=6)
add_image(s, OUT / "phase3/fig_enrichment_coverage.png", 6.5, 1.4, width=6.3)
caption_box(s, 0.3, 5.8, 12.5,
            "World Bank (GDP, internet, mobile, education), UN HDI, LDS Church Statistics, "
            "and Pew Research (religious composition, restrictions, behavioral religiosity) provide contextual features "
            "for 99%+ of users. These form the H0 feature block for discriminant analysis.")
print("Slide 5: Enrichment")

# ═══════════════════════════════════════════════════
# SLIDE 6: GDP/HDI/LDS Landscape
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "The Global Development Landscape of FamilySearch Users")
add_image(s, OUT / "phase3/fig_gdp_hdi_lds.png", 1.5, 1.4, width=10)
caption_box(s, 1.5, 6.0, 10,
            "Each point is a country. Color intensity = LDS member density (log scale). "
            "FamilySearch users span the full spectrum from low-HDI/low-GDP (bottom-left) to high-HDI/high-GDP (top-right). "
            "LDS density is highest in middle-income Latin American countries and the US.")
print("Slide 6: GDP/HDI/LDS")

# ═══════════════════════════════════════════════════
# SLIDE 7: Persistence — The Target Variable
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Defining Persistence: A Composite of Login Consistency, Recency, and Activity Breadth")
add_image(s, OUT / "phase2/fig_persistence_c_dist.png", 1.5, 1.4, width=10)
caption_box(s, 1.5, 5.8, 10,
            "Persistence Score C combines login frequency (how often), recency of last milestone (how recently), "
            "and activity breadth (how many types). The distribution is right-skewed: most users cluster near the floor "
            "(low persistence), with a long tail of highly persistent contributors. "
            "Median split within the analytical population defines Persistent vs Transient for classification.")
print("Slide 7: Persistence Distribution")

# ═══════════════════════════════════════════════════
# SLIDE 8: PCA — Three Orthogonal Axes (All Tier D)
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "The Data Has Three Independent Dimensions: Engagement, Timing, and Context",
          "All Tier D users (including 1-login) — PC1: Volume, PC2: Velocity, PC3: Context")
add_image(s, OUT / "exploratory/fig_pca_loadings_pc1.png", 0.1, 1.3, width=4.3)
add_image(s, OUT / "exploratory/fig_pca_loadings_pc2.png", 4.5, 1.3, width=4.3)
add_image(s, OUT / "exploratory/fig_pca_loadings_pc3.png", 8.9, 1.3, width=4.3)
caption_box(s, 0.3, 5.8, 12.5,
            "PC1 (20.4%): Engagement Volume — all top features are behavioral contribution rates. "
            "PC2 (13.0%): Onboarding Velocity — milestone timing features. "
            "PC3 (11.7%): Country Context — GDP, HDI, GEPI, religious diversity. "
            "These axes are nearly orthogonal: behavioral engagement and country context occupy different dimensions of the data.")
print("Slide 8: Three Axes")

# ═══════════════════════════════════════════════════
# SLIDE 9: Biplot — Features Map to Persistence Direction
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Volume Features Point Toward Persistence; Context Features Are Orthogonal",
          "All Tier D — PCA Biplot with Feature Loading Vectors")
add_image(s, OUT / "exploratory/fig_biplot.png", 1.5, 1.3, width=10)
caption_box(s, 1.5, 6.2, 10,
            "Green arrows (behavioral Volume features) align with the persistence gradient (green=persistent users). "
            "Blue arrows (enrichment/contextual) point into a perpendicular dimension — they explain different variance than persistence. "
            "This geometric separation is why contextual features add nothing to persistence prediction.")
print("Slide 9: Biplot")

# ═══════════════════════════════════════════════════
# SLIDE 10: H1 vs H0 — First Pass
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Behavioral Features Classify Persistence Near-Perfectly; Context Is No Better Than Chance",
          "Phase 5: All Tier D — 6 Feature Blocks × 3 Models × 10 Subsamples = 180 Model Runs")
add_image(s, OUT / "phase5/fig_block_auc_comparison.png", 0.3, 1.3, width=12.5)
caption_box(s, 0.3, 6.0, 12.5,
            "Block 4 (H1: Velocity+Volume+Sequencing) AUC = 0.999. Block 5 (H0: Age, Country, GDP, HDI, Religiosity) AUC = 0.59. "
            "Block 6 (Full model) = 0.999 — adding contextual features to engagement provides zero lift. "
            "Delta_H1 = +0.41 (engagement adds massively to context). Delta_H0 = 0.00 (context adds nothing to engagement).")
print("Slide 10: H1 vs H0")

# ═══════════════════════════════════════════════════
# SLIDE 11: First-Pass Clustering
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Unsupervised Clusters Naturally Separate Persistent from Transient Users",
          "Phase 6: K-Means k=3 on behavioral features — All Tier D")
add_image(s, OUT / "phase6/fig_pca_clusters.png", 0.1, 1.3, width=6.4)
add_image(s, OUT / "phase6/fig_pca_persistence.png", 6.6, 1.3, width=6.4)
caption_box(s, 0.3, 5.8, 12.5,
            "Left: 3 clusters discovered by K-Means on behavioral features. "
            "Right: Same users colored by persistence score (red=low, green=high). "
            "The cluster boundaries align with the persistence gradient — Cramer's V = 0.455 (strong association). "
            "Clusters: Minimal Engagers (55%, 64% persistent), Mid-Range (28%, 82%), Power Contributors (17%, 94%).")
print("Slide 11: Clustering")

# ═══════════════════════════════════════════════════
# SLIDE 12: Contributors Only — Refined Classification
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Removing Single-Visit Users Strengthens H1 and Reveals Richer Structure",
          "Phase 5b: Contributors Only (2+ logins) — a harder classification task")
add_image(s, OUT / "phase5b/fig_block_auc_comparison.png", 0.3, 1.3, width=12.5)
caption_box(s, 0.3, 6.0, 12.5,
            "After removing 51% of Tier D users who logged in only once, the classification task is harder (higher persistence floor). "
            "Yet H1 features still achieve AUC 0.997. H0 drops to 0.538 — even closer to random chance. "
            "Sequencing gains +0.067 AUC (from 0.69 to 0.76); Velocity features enter the top 10 for the first time. "
            "The 1-login population was masking real Sequencing and Velocity signal.")
print("Slide 12: Contributors Classification")

# ═══════════════════════════════════════════════════
# SLIDE 13: PC2/PC3 Axis Swap
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Removing Single-Visit Users Reveals That Context Is the Second-Largest Data Dimension",
          "Contributors Only: PC2 becomes Contextual, PC3 becomes Velocity")
add_image(s, OUT / "exploratory_b/fig_pca_loadings_pc2.png", 0.1, 1.3, width=6.4)
add_image(s, OUT / "exploratory_b/fig_pca_loadings_pc3.png", 6.6, 1.3, width=6.4)
caption_box(s, 0.3, 5.8, 12.5,
            "In the all-Tier-D analysis, PC2 was Velocity and PC3 was Context. After removing 1-login users, they swap. "
            "PC2 is now dominated by enrichment features (GEPI 0.46, GDP 0.46, HDI 0.36). "
            "PC3 becomes Velocity (days_to_first_name 0.51). The 1-login population inflated Velocity variance; "
            "without it, country development becomes the second-largest independent dimension of user variation.")
print("Slide 13: Axis Swap")

# ═══════════════════════════════════════════════════
# SLIDE 14: Contributors Biplot + Country Overlap
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Country Clusters Overlap in Behavioral Space — Geography Doesn't Partition Engagement",
          "Contributors Only — PCA projections")
add_image(s, OUT / "exploratory_b/fig_biplot.png", 0.1, 1.3, width=6.4)
add_image(s, OUT / "exploratory_b/fig_pca_full_country.png", 6.6, 1.3, width=6.4)
caption_box(s, 0.3, 5.8, 12.5,
            "Left: Biplot — Volume arrows point toward persistence; enrichment arrows are orthogonal. "
            "Right: Same projection colored by country cluster. All clusters overlap extensively — "
            "users from 'High-LDS International' and 'Mod-Eng Low-LDS' occupy the same behavioral feature space. "
            "No geographic separation exists in the engagement dimensions.")
print("Slide 14: Biplot + Country")

# ═══════════════════════════════════════════════════
# SLIDE 15: 6-Cluster Solution
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Contributors Reveal 6 Behavioral Segments with 54-100% Persistence Rates",
          "Phase 6b: K-Means k=6 — Cramer's V = 0.574 (very strong cluster↔persistence association)")
add_image(s, OUT / "phase6b/fig_pca_clusters.png", 0.1, 1.3, width=6.4)
add_image(s, OUT / "phase6b/fig_pca_persistence.png", 6.6, 1.3, width=6.4)
caption_box(s, 0.3, 5.8, 12.5,
            "Left: 6 clusters on behavioral features. Right: Persistence gradient. "
            "The 1-login removal revealed 3 sub-segments hidden in the first-pass 'Minimal Engagers' cluster. "
            "Cluster 3 (n=159) has 99.6% persistence with only moderate logins — an ultra-stable core worth understanding.")
print("Slide 15: 6 Clusters")

# ═══════════════════════════════════════════════════
# SLIDE 16: Radar Profiles
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Six Distinct Behavioral Profiles — From Light Engagers to Power Contributors")
add_image(s, OUT / "phase6b/fig_radar_profiles.png", 2.5, 1.3, width=8)
caption_box(s, 2.0, 5.8, 9,
            "Radar charts show normalized feature means per cluster. Power Contributors (Cluster 5) score highest on breadth and sources. "
            "The 'efficient persisters' (Cluster 4) achieve 94% persistence with moderate volume — efficiency over intensity. "
            "Cluster 0 (Light Engagers, 54% persistent) is the primary churn risk and the largest segment (45%).")
print("Slide 16: Radar")

# ═══════════════════════════════════════════════════
# SLIDE 17: Development Tiers × Persistence Gradient — THE KEY FINDING
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "The Engagement→Persistence Relationship Is Strongest in Middle-Development Countries",
          "Segmentation: 5 contextual-development tiers × Volume gradient")
add_image(s, OUT / "segmentation/fig_gradient_by_tier.png", 0.3, 1.3, width=8.5)
add_text(s, 9.0, 1.5, 4.0, 5.2,
         "Tiers defined by K-Means on PC2\n(Contextual axis: GDP, HDI, GEPI)\n\n"
         "T1 (High-Dev): slope 0.012\n"
         "Users persist regardless of volume.\n"
         "Engagement is confirmation,\nnot cause.\n\n"
         "T3 (Mid-Dev): slope 0.038\n"
         "Strongest gradient — persistence\nis most sensitive to engagement\nvolume. The intervention tier.\n\n"
         "T5 (Low-Dev): slope 0.031\n"
         "High variance — widest persistence\nrange (0.73). Some become power\nusers; many churn.\n\n"
         "Interaction: F=190, p≈0\n"
         "R² gain: +2.3%",
         font_size=12, color=BLACK)
print("Slide 17: Gradient by Tier")

# ═══════════════════════════════════════════════════
# SLIDE 18: Conclusions
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s, GREEN_DARK)
add_text(s, 0.8, 0.5, 12, 0.8, "Conclusions", font_size=36, bold=True, color=WHITE)
add_text(s, 0.8, 1.4, 11.5, 5.5,
         "1. Behavioral engagement patterns predict user Persistence with near-perfect accuracy (AUC 0.999). "
         "Volume — specifically login frequency and contribution rate in the first 90 days — accounts for 82% of "
         "discriminant importance. Velocity and Sequencing contribute 14% and 5% respectively among multi-session contributors.\n\n"
         "2. Demographic and contextual factors (age, country, GDP, HDI, religiosity) achieve AUC 0.54-0.63 — "
         "barely above chance. Adding them to the behavioral model provides zero incremental lift.\n\n"
         "3. However, contextual factors DO naturally segment the user population into distinct development tiers. "
         "Within these tiers, the engagement→persistence relationship is strongly linear but varies in strength: "
         "T1 (high-development) shows a weak, plateauing gradient — users persist regardless of engagement volume. "
         "T3 (middle-development) shows the strongest gradient — these users are most responsive to engagement interventions.\n\n"
         "4. The finding suggests that FamilySearch's retention strategy should be context-adaptive: "
         "maximize early engagement volume in middle-development markets (where the marginal return is highest) "
         "while empowering already-committed users in high-development markets with contributor tools and community features.\n\n"
         "5. Lower-tier segments have fewer 'power contributors' — suggesting untapped potential that could benefit from "
         "structured introductions to FamilySearch's deeper features: source attachment, memory creation, and Get Involved campaigns.",
         font_size=14, color=WHITE)
print("Slide 18: Conclusions")

# ═══════════════════════════════════════════════════
# SLIDE 19: Future Work & Data Wishes
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Future Work & The Data We Wish We Had")
add_text(s, 0.5, 1.5, 6.0, 5.5,
         "Future Work\n\n"
         "• Longitudinal validation: track cohorts over 12+ months to confirm "
         "that 90-day behavioral patterns predict actual 1-year retention\n\n"
         "• Causal inference: A/B test engagement nudges in T3 (middle-development) "
         "tier to measure causal impact on persistence\n\n"
         "• Non-login contributor analysis: the 5.7% of users who contribute "
         "without logging in suggest an API/batch user journey worth understanding\n\n"
         "• Survival analysis: model time-to-churn as a continuous outcome "
         "rather than binary Persistent/Transient\n\n"
         "• ACCOUNT_TYPE cross-validation: assess whether the discriminant "
         "function works equally well for Member vs Public accounts\n\n"
         "• Production-scale validation: T=5 × 30K subsamples with full "
         "R=200 bootstrap stability testing",
         font_size=13, color=BLACK)
add_text(s, 6.8, 1.5, 6.0, 5.5,
         "The Data We Wish We Had\n\n"
         "• Per-session activity logs (not just cumulative counts and "
         "first-occurrence dates) — would enable true longitudinal trajectory "
         "analysis and time-series clustering\n\n"
         "• Last activity date (not just first) — Persistence Definition C "
         "uses MAX(EARLIEST_*) as a recency proxy, which may underestimate "
         "users who are still active on already-established activity types\n\n"
         "• Device type (mobile vs desktop) — the literature predicts "
         "shorter sessions and higher frequency for mobile-first users\n\n"
         "• Referral source / acquisition channel — church-directed "
         "sign-ups likely have different persistence profiles than organic\n\n"
         "• Gender — a potential confound for engagement analysis, "
         "completely absent from this dataset\n\n"
         "• Google Trends genealogy search volume by country — would "
         "complete the GEPI composite (deferred: API access blocked)",
         font_size=13, color=BLACK)
print("Slide 19: Future Work")

# ═══════════════════════════════════════════════════
# SLIDE 20: Scope & Constraints
# ═══════════════════════════════════════════════════
s = add_slide()
add_bg(s)
title_bar(s, "Scope Constraints & Methodological Notes")
add_text(s, 0.5, 1.5, 12, 5.5,
         "Proof-of-Concept Scope\n"
         "• Subsamples: T=10 × 5,079 users (production would use T=5 × 30,000)\n"
         "• Bootstrap stability: Jaccard via cross-subsample matching (production: R=200 clusterboot iterations)\n"
         "• Persistence score: composite proxy using EARLIEST dates only, not true survival curve\n\n"
         "Data Limitations\n"
         "• MNAR block (10.2%): excluded, not imputed — these users may differ systematically from tracked users\n"
         "• Cross-sectional data only: no per-period activity counts; Velocity construct limited to 3 milestone transitions\n"
         "• Province/City data exists exclusively for Member accounts (3%) — excluded from modeling to avoid account-type leakage\n"
         "• Pew behavioral religiosity covers only 24 countries (53% of users); composition/restrictions cover 99%+\n"
         "• Fixed 90-day Volume window is approximate (prorated from cumulative counts, not true windowed measurement)\n\n"
         "Analytical Decisions\n"
         "• ACCOUNT_TYPE excluded from discriminant models to avoid trivial Member/Public classification\n"
         "• Persistence dichotomized within Tier D population (not full population) to avoid floor effects\n"
         "• VIF > 10 for 13 features (constructed collinearity); does not affect block-level AUC comparisons\n"
         "• HDBSCAN deferred: K-Means and GMM agreed on optimal k; gradient structure favors continuous methods\n\n"
         "Reproducibility\n"
         "• All random operations seeded and recorded in experiment_registry table\n"
         "• All intermediate results exported as Parquet/CSV; DuckDB DDL tracked in git\n"
         "• Phase scripts executable independently: src/phase{1-6}*.py",
         font_size=12, color=BLACK)
print("Slide 20: Scope")

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
prs.save(str(DECK_PATH))
print(f"\n=== Deck saved: {DECK_PATH} ===")
print(f"Slides: {len(prs.slides)}")
